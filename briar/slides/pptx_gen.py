"""Briar Slide Generator — PowerPoint presentations 🎨"""
import os
from datetime import datetime

class SlideGenerator:
    def __init__(self, target: str, findings: list, output_dir: str = "./reports/slides"):
        self.target = target
        self.findings = findings
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def pptx(self) -> str:
        """Generate PowerPoint presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Slide 1 — Title
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Security Audit Report"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(184, 71, 240)
            p.alignment = PP_ALIGN.CENTER
            
            p2 = tf.add_paragraph()
            p2.text = self.target
            p2.font.size = Pt(28)
            p2.font.color.rgb = RGBColor(205, 214, 244)
            p2.alignment = PP_ALIGN.CENTER
            
            p3 = tf.add_paragraph()
            p3.text = f"Briar v0.1.0 — {datetime.now().strftime('%Y-%m-%d')}"
            p3.font.size = Pt(18)
            p3.font.color.rgb = RGBColor(108, 112, 134)
            p3.alignment = PP_ALIGN.CENTER
            
            # Slide 2 — Summary
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
            tf2 = txBox2.text_frame
            tf2.paragraphs[0].text = "Executive Summary"
            tf2.paragraphs[0].font.size = Pt(36)
            tf2.paragraphs[0].font.bold = True
            tf2.paragraphs[0].font.color.rgb = RGBColor(184, 71, 240)
            
            items = [
                f"Target: {self.target}",
                f"Total findings: {len(self.findings)}",
                f"Scan date: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            ]
            for item in items:
                p = tf2.add_paragraph()
                p.text = f"  {item}"
                p.font.size = Pt(20)
                p.space_after = Pt(12)
            
            # Slides for each critical finding
            for i, finding in enumerate(self.findings[:5]):
                if finding.get("severity") in ["Critical", "High"]:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(6))
                    tf = txBox.text_frame
                    tf.paragraphs[0].text = f"Finding {i+1}: {finding.get('type','Vulnerability')}"
                    tf.paragraphs[0].font.size = Pt(28)
                    tf.paragraphs[0].font.bold = True
                    
                    p = tf.add_paragraph()
                    p.text = f"Severity: {finding.get('severity','Unknown')}"
                    p.font.size = Pt(18)
                    
                    analysis = finding.get('analysis', finding.get('raw',''))[:500]
                    p2 = tf.add_paragraph()
                    p2.text = analysis
                    p2.font.size = Pt(14)
            
            # Last slide — Recommendations
            slide_end = prs.slides.add_slide(prs.slide_layouts[6])
            txBox_end = slide_end.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
            tf_end = txBox_end.text_frame
            tf_end.paragraphs[0].text = "Recommendations"
            tf_end.paragraphs[0].font.size = Pt(36)
            tf_end.paragraphs[0].font.bold = True
            tf_end.paragraphs[0].font.color.rgb = RGBColor(166, 227, 161)
            
            recs = [
                "Review all findings and prioritize by severity",
                "Apply remediation steps for each vulnerability",
                "Re-scan after fixes are applied",
                "Schedule regular automated scans with Briar",
            ]
            for r in recs:
                p = tf_end.add_paragraph()
                p.text = f"  {r}"
                p.font.size = Pt(20)
                p.space_after = Pt(16)
            
            path = os.path.join(self.output_dir, f"slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
            prs.save(path)
            return path
        except ImportError:
            return "pip install python-pptx"
    
    def html_slides(self) -> str:
        """Generate HTML/Canva-like slides"""
        html = f"""<!DOCTYPE html>
<html lang="en"><head><meta charset="UTF-8"><title>Briar Report — {self.target}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0a0514;color:#cdd6f4;font-family:'JetBrains Mono',monospace}}
.slide{{min-height:100vh;display:flex;align-items:center;justify-content:center;padding:3rem;border-bottom:1px solid #1e1e2e}}
.slide-content{{max-width:900px;width:100%}}
h1{{color:#B847F0;font-size:3rem;margin-bottom:1rem}}
h2{{color:#B847F0;font-size:2rem;margin-bottom:1rem}}
.finding{{background:#1e1e2e;border:1px solid #313244;border-radius:12px;padding:1.5rem;margin:1rem 0}}
.critical{{border-color:#f38ba8}}.high{{border-color:#fab387}}
.meta{{color:#6c7086;font-size:0.9rem}}
</style></head><body>
<div class="slide"><div class="slide-content">
<h1>Security Audit Report</h1>
<p style="font-size:1.5rem">{self.target}</p>
<p class="meta">Briar v0.1.0 — {datetime.now().strftime('%Y-%m-%d')}</p>
<p class="meta">Findings: {len(self.findings)}</p>
</div></div>
"""
        for i, f in enumerate(self.findings):
            sev = f.get("severity","").lower()
            html += f"""
<div class="slide"><div class="slide-content">
<h2>Finding {i+1}: {f.get('type','?')}</h2>
<div class="finding {sev}">
<p><strong>Severity:</strong> {f.get('severity','?')}</p>
<p><strong>Agent:</strong> {f.get('agent','?')}</p>
<p>{f.get('analysis', f.get('raw',''))[:300]}...</p>
</div></div></div>"""
        html += "</body></html>"
        
        path = os.path.join(self.output_dir, f"slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html")
        with open(path, "w") as f:
            f.write(html)
        return path
